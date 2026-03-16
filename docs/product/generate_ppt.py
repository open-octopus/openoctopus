#!/usr/bin/env python3
"""Generate OpenOctopus Business Plan PPT."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === Brand Colors ===
BG_DARK = RGBColor(0x0A, 0x19, 0x29)
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x88, 0x99, 0xAA)
GREEN = RGBColor(0x00, 0xC8, 0x53)
RED = RGBColor(0xFF, 0x45, 0x45)
YELLOW = RGBColor(0xFF, 0xD6, 0x00)
LIGHT_BG = RGBColor(0x0D, 0x1F, 0x33)
CARD_BG = RGBColor(0x12, 0x2A, 0x40)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = spacing
    return txBox


def add_card(slide, left, top, width, height, title, items, title_color=ACCENT_CYAN, bg_color=CARD_BG):
    add_shape(slide, left, top, width, height, fill_color=bg_color)
    add_text(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
             title, font_size=16, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.4), height - Inches(0.7),
                    items, font_size=13, color=WHITE, spacing=Pt(4))


def slide_title(slide, title, subtitle=None):
    add_shape(slide, Inches(0), Inches(0), W, Inches(1.2), fill_color=LIGHT_BG)
    add_text(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.6), title, font_size=32, color=ACCENT_CYAN, bold=True)
    if subtitle:
        add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.4), subtitle, font_size=16, color=GRAY)


def add_table(slide, left, top, width, row_height, headers, rows, col_widths=None):
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, row_height * num_rows)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = "Arial"
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x15, 0x35, 0x55)

    # Rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = WHITE
                p.font.name = "Arial"
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if r % 2 == 0 else LIGHT_BG

    return table_shape


# ============================================================
# Slide 1: Cover
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# Decorative top bar
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), fill_color=ACCENT_CYAN)

# Title
add_text(slide, Inches(0), Inches(1.8), W, Inches(1.0),
         "OpenOctopus", font_size=60, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(2.8), W, Inches(0.6),
         "商 业 计 划 书", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(3.6), W, Inches(0.5),
         "全球首个 Realm-native 人生治理 Agent 系统", font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

# Decorative line
add_shape(slide, Inches(5.0), Inches(4.3), Inches(3.333), Inches(0.04), fill_color=ACCENT_CYAN)

# Bottom info
add_text(slide, Inches(0), Inches(5.8), W, Inches(0.4),
         "2026年3月", font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(6.2), W, Inches(0.4),
         "hello@openoctopus.org  |  github.com/open-octopus", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# Bottom bar
add_shape(slide, Inches(0), Inches(7.42), W, Inches(0.08), fill_color=ACCENT_CYAN)


# ============================================================
# Slide 2: Executive Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "执行摘要", "Executive Summary")

# Core statement
add_shape(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.9), fill_color=CARD_BG, line_color=ACCENT_CYAN)
add_text(slide, Inches(1.0), Inches(1.5), Inches(11.3), Inches(0.7),
         "将你的人生分域治理（宠物、父母、财务、工作...），通过 Summon 将任何现实对象转化为有记忆、有性格、可主动行动的 AI Agent。",
         font_size=18, color=WHITE)

# Three pain-point / solution cards
cards_data = [
    ("AI工具碎片化", ["❌ ChatGPT/Notion/记账软件来回切换", "✅ Realm Matrix 一站式管理"]),
    ("情感连接缺失", ["❌ 通用AI冷冰冰", "✅ Summon 实体人格化"]),
    ("隐私焦虑", ["❌ 云端AI存储敏感数据", "✅ Local-first 本地存储"]),
]
for i, (title, items) in enumerate(cards_data):
    x = Inches(0.8 + i * 4.0)
    add_card(slide, x, Inches(2.6), Inches(3.6), Inches(1.6), title, items, title_color=ACCENT_ORANGE)

# Key metrics
metrics = [
    ("🎯 目标市场", "$7.6B → $183B (2033)"),
    ("💰 融资需求", "种子轮 ¥200万"),
    ("👥 团队规模", "5人核心团队"),
    ("📊 开发进度", "Phase 1 完成"),
]
for i, (label, value) in enumerate(metrics):
    x = Inches(0.8 + i * 3.1)
    add_shape(slide, x, Inches(4.6), Inches(2.8), Inches(1.3), fill_color=CARD_BG)
    add_text(slide, x, Inches(4.7), Inches(2.8), Inches(0.4), label, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(5.15), Inches(2.8), Inches(0.5), value, font_size=18, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 3: Summon Mechanism
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "核心创新 — Summon 机制", "将现实实体转化为 AI Agent")

# Flow chart
flow_steps = ["原始数据", "结构化实体", "Summon", "有记忆·有个性·可行动的Agent"]
for i, step in enumerate(flow_steps):
    x = Inches(0.8 + i * 3.15)
    color = ACCENT_CYAN if step == "Summon" else CARD_BG
    line = ACCENT_CYAN if step == "Summon" else RGBColor(0x30, 0x50, 0x70)
    add_shape(slide, x, Inches(1.5), Inches(2.8), Inches(0.7), fill_color=color, line_color=line)
    fc = BG_DARK if step == "Summon" else WHITE
    add_text(slide, x, Inches(1.55), Inches(2.8), Inches(0.6), step, font_size=16, color=fc, bold=True, alignment=PP_ALIGN.CENTER)
    if i < 3:
        add_text(slide, Inches(3.5 + i * 3.15), Inches(1.55), Inches(0.5), Inches(0.6), "→", font_size=24, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# 4 entity types (2x2)
entity_types = [
    ("🐕 Living", "宠物、家人、朋友", "模拟对话、情感表达"),
    ("🏠 Asset", "汽车、房产、投资组合", "状态监控、维护提醒"),
    ("🏢 Organization", "医院、学校、公司", "流程向导、预约协调"),
    ("📋 Abstract", "目标、项目、习惯", "进度追踪、偏差预警"),
]
for i, (title, desc, ability) in enumerate(entity_types):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 3.3)
    y = Inches(2.6 + row * 1.5)
    add_card(slide, x, y, Inches(3.0), Inches(1.3), title, [desc, ability])

# Dialog example
add_shape(slide, Inches(7.0), Inches(2.6), Inches(5.5), Inches(4.2), fill_color=CARD_BG, line_color=ACCENT_CYAN)
add_text(slide, Inches(7.2), Inches(2.7), Inches(5.1), Inches(0.4),
         '用户："我下周要出差5天。"', font_size=15, color=WHITE, bold=True)

dialogs = [
    ('🐕 Momo', '"谁要喂我、遛我5天啊？！"', ACCENT_CYAN),
    ('👩 Mom', '"妈说她正好想来看看——她可以照顾 Momo。"', GREEN),
    ('🚗 Car', '"出发前充满电？还是打车去机场？"', YELLOW),
    ('💰 Budget', '"预估差旅费 ¥3,500，记得保存发票。"', ACCENT_ORANGE),
]
for i, (speaker, text, color) in enumerate(dialogs):
    y = Inches(3.3 + i * 0.85)
    add_text(slide, Inches(7.4), y, Inches(2.0), Inches(0.35), speaker, font_size=14, color=color, bold=True)
    add_text(slide, Inches(7.4), y + Inches(0.3), Inches(5.0), Inches(0.4), text, font_size=13, color=WHITE)


# ============================================================
# Slide 4: OpenClaw Comparison
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "与 OpenClaw 对比", "OpenClaw 已验证市场，但存在空白")

# OpenClaw (left)
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.5), fill_color=CARD_BG, line_color=RED)
add_text(slide, Inches(1.0), Inches(1.6), Inches(5.1), Inches(0.4),
         "OpenClaw 现状", font_size=20, color=RED, bold=True)
openclaw_items = [
    '✅ 验证了"本地优先 + Agent"需求',
    "❌ 技术门槛过高（需编写复杂配置）",
    "❌ 假冒仓库传播恶意软件",
    "❌ 定位偏向任务自动化",
]
add_bullet_list(slide, Inches(1.0), Inches(2.2), Inches(5.1), Inches(2.5), openclaw_items, font_size=16)

# OpenOctopus (right)
add_shape(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(3.5), fill_color=CARD_BG, line_color=GREEN)
add_text(slide, Inches(7.2), Inches(1.6), Inches(5.1), Inches(0.4),
         "OpenOctopus 填补空白", font_size=20, color=GREEN, bold=True)
oo_items = [
    "✅ 零门槛 Summon 机制",
    "✅ 官方正版，安全审核",
    '✅ 聚焦"人生治理"场景',
    "✅ 扩大10倍潜在用户群",
]
add_bullet_list(slide, Inches(7.2), Inches(2.2), Inches(5.1), Inches(2.5), oo_items, font_size=16)

# Comparison table
add_table(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.5),
          ["维度", "OpenClaw", "OpenOctopus"],
          [
              ["定位", "任务自动化工具", "人生治理系统"],
              ["用户门槛", "高（写配置）", "低（一键召唤）"],
              ["目标用户", "技术极客", "生活管理者"],
          ],
          col_widths=[Inches(2.5), Inches(4.6), Inches(4.6)])


# ============================================================
# Slide 5: Technical Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "技术架构", "Realm-native 架构")

# Central Brain
add_shape(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(0.7), fill_color=ACCENT_CYAN)
add_text(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(0.7),
         "Central Brain", font_size=18, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)

# Sub-components
sub_comps = ["Router", "Coordinator", "Scheduler"]
for i, name in enumerate(sub_comps):
    x = Inches(3.5 + i * 2.3)
    add_shape(slide, x, Inches(2.5), Inches(2.0), Inches(0.55), fill_color=LIGHT_BG, line_color=ACCENT_CYAN)
    add_text(slide, x, Inches(2.5), Inches(2.0), Inches(0.55), name, font_size=14, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# Realm boxes
realms = ["Pet Realm", "Finance Realm", "Legal Realm", "Parents Realm"]
realm_icons = ["🐕", "💰", "⚖️", "👨‍👩‍👧"]
for i, (name, icon) in enumerate(zip(realms, realm_icons)):
    x = Inches(0.8 + i * 3.15)
    add_shape(slide, x, Inches(3.5), Inches(2.8), Inches(2.0), fill_color=CARD_BG, line_color=ACCENT_CYAN)
    add_text(slide, x, Inches(3.55), Inches(2.8), Inches(0.5),
             f"{icon} {name}", font_size=16, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    add_bullet_list(slide, x + Inches(0.3), Inches(4.1), Inches(2.2), Inches(1.2),
                    ["Agents", "Skills", "Memory"], font_size=13, color=GRAY, spacing=Pt(4))

# 12 Realms icons
add_text(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.4),
         "12个预设 Realm", font_size=16, color=WHITE, bold=True)
realms_12 = "🐕 Pet  |  👨‍👩‍👧 Parents  |  💕 Partner  |  💰 Finance  |  💼 Work  |  ⚖️ Legal  |  🚗 Vehicle  |  🏠 Home  |  🏥 Health  |  🏃 Fitness  |  🎨 Hobby  |  👥 Friends"
add_text(slide, Inches(0.8), Inches(6.15), Inches(11.7), Inches(0.4), realms_12, font_size=14, color=GRAY)

# Tech stack
tech = "Node.js ≥22 + TypeScript  |  SQLite (Local-first)  |  Express 5 + WebSocket  |  MIT 开源协议"
add_shape(slide, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5), fill_color=LIGHT_BG)
add_text(slide, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5), tech, font_size=13, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 6: Market Analysis
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "市场分析", "AI Agent 市场高速增长")

# Big numbers
add_text(slide, Inches(0.8), Inches(1.5), Inches(4.0), Inches(1.0),
         "$7.6B", font_size=54, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(2.5), Inches(2.0), Inches(0.4),
         "2025年", font_size=18, color=GRAY)

add_text(slide, Inches(4.5), Inches(1.8), Inches(1.5), Inches(0.6),
         "→", font_size=40, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(5.5), Inches(1.5), Inches(4.5), Inches(1.0),
         "$183B", font_size=54, color=ACCENT_CYAN, bold=True)
add_text(slide, Inches(5.5), Inches(2.5), Inches(3.0), Inches(0.4),
         "2033年  CAGR 49.6%", font_size=18, color=GRAY)

add_text(slide, Inches(0.8), Inches(3.0), Inches(6.0), Inches(0.3),
         "数据来源：Grand View Research, 2025", font_size=12, color=GRAY)

# TAM/SAM/SOM pyramid (using stacked rectangles)
pyramid_data = [
    ("SOM: $50M (2026)", "早期采用者", Inches(4.8), Inches(1.6), ACCENT_CYAN),
    ("SAM: $500M (2026)", "个人AI助手细分市场", Inches(3.8), Inches(2.6), RGBColor(0x00, 0x90, 0xBB)),
    ("TAM: $183B (2033)", "全球AI Agent市场", Inches(2.8), Inches(3.6), RGBColor(0x00, 0x60, 0x88)),
]
for label, desc, x, w, color in pyramid_data:
    cx = x + (Inches(6.6) - w) / 2 + Inches(3.5)
    # calculate y position
    idx = pyramid_data.index((label, desc, x, w, color))
    y = Inches(3.6 + idx * 1.1)
    add_shape(slide, cx, y, w, Inches(0.9), fill_color=color)
    add_text(slide, cx, y + Inches(0.05), w, Inches(0.45), label, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, cx, y + Inches(0.45), w, Inches(0.4), desc, font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)

# User personas
add_text(slide, Inches(0.8), Inches(3.5), Inches(5.0), Inches(0.4),
         "目标用户画像", font_size=18, color=WHITE, bold=True)

# Persona 1
add_card(slide, Inches(0.8), Inches(4.0), Inches(5.0), Inches(1.5),
         "核心用户：极客/技术爱好者",
         ["25-40岁，开发者/产品经理", "使用多个AI工具，感到碎片化", "关注隐私，愿为个性化付费"],
         title_color=ACCENT_CYAN)

# Persona 2
add_card(slide, Inches(0.8), Inches(5.7), Inches(5.0), Inches(1.5),
         "扩展用户：生活管理者",
         ["30-50岁，多角色身份", "管理复杂生活：宠物、父母、投资", "非技术背景，需要低门槛"],
         title_color=ACCENT_ORANGE)


# ============================================================
# Slide 7: Competitive Analysis
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "竞品矩阵分析", "Competitive Analysis")

# Positioning chart (simplified with shapes)
# Axes
chart_left = Inches(0.8)
chart_top = Inches(1.5)
chart_w = Inches(5.5)
chart_h = Inches(4.5)
add_shape(slide, chart_left, chart_top, chart_w, chart_h, fill_color=LIGHT_BG)

# Labels
add_text(slide, chart_left - Inches(0.1), chart_top - Inches(0.05), Inches(1.5), Inches(0.3), "高技术门槛", font_size=11, color=GRAY)
add_text(slide, chart_left - Inches(0.1), chart_top + chart_h - Inches(0.3), Inches(1.5), Inches(0.3), "低技术门槛", font_size=11, color=GRAY)
add_text(slide, chart_left + Inches(0.1), chart_top + chart_h + Inches(0.05), Inches(1.5), Inches(0.3), "个人效率", font_size=11, color=GRAY)
add_text(slide, chart_left + chart_w - Inches(1.2), chart_top + chart_h + Inches(0.05), Inches(1.5), Inches(0.3), "人生治理", font_size=11, color=GRAY)

# Competitor dots
competitors = [
    ("OpenClaw", Inches(2.5), Inches(0.8), RED),
    ("Notion AI", Inches(1.2), Inches(2.8), GRAY),
    ("Coze", Inches(3.5), Inches(3.2), GRAY),
    ("OpenOctopus ★", Inches(4.0), Inches(1.5), ACCENT_CYAN),
]
for name, rx, ry, color in competitors:
    cx = chart_left + rx
    cy = chart_top + ry
    dot = add_shape(slide, cx, cy, Inches(0.25), Inches(0.25), fill_color=color, shape_type=MSO_SHAPE.OVAL)
    add_text(slide, cx + Inches(0.3), cy - Inches(0.05), Inches(1.8), Inches(0.3), name, font_size=12, color=color, bold=True)

# Comparison table (right side)
add_table(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(0.5),
          ["产品", "类型", "核心差异"],
          [
              ["OpenClaw", "开源Agent", "任务导向，配置复杂"],
              ["Notion AI", "工作空间", "文档为中心，无领域划分"],
              ["Sintra AI", "商业Agent", "仅面向商业"],
              ["Coze/扣子", "Agent平台", "无领域组织，无本地优先"],
          ],
          col_widths=[Inches(1.4), Inches(1.5), Inches(2.8)])

# Moat
add_text(slide, Inches(6.8), Inches(4.2), Inches(5.5), Inches(0.4),
         "我们的护城河", font_size=18, color=WHITE, bold=True)
moat_items = [
    "🏰 Realm-native 架构专利潜力",
    "🤖 Summon 实体转化引擎",
    "🕸️ 跨域知识图谱",
    "🔒 本地优先 + 云同步混合架构",
]
add_bullet_list(slide, Inches(6.8), Inches(4.7), Inches(5.5), Inches(2.0), moat_items, font_size=15, color=ACCENT_CYAN)


# ============================================================
# Slide 8: Business Model
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "商业模式", "开源 + 增值服务")

# Pricing tiers
tiers = [
    ("社区版", "¥0", "开源核心功能", GRAY),
    ("Pro版", "¥29/月", "高级Agent + 云同步", ACCENT_CYAN),
    ("家庭版", "¥69/月", "多用户 + 家庭共享", GREEN),
    ("企业版", "¥199/人/月", "企业部署 + 专属支持", ACCENT_ORANGE),
]
for i, (name, price, desc, color) in enumerate(tiers):
    x = Inches(0.8 + i * 3.1)
    h = Inches(2.2)
    add_shape(slide, x, Inches(1.5), Inches(2.8), h, fill_color=CARD_BG, line_color=color)
    add_text(slide, x, Inches(1.6), Inches(2.8), Inches(0.4), name, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(2.1), Inches(2.8), Inches(0.5), price, font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(2.7), Inches(2.8), Inches(0.4), desc, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# Revenue breakdown
add_text(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(0.4),
         "2027年收入占比预测", font_size=18, color=WHITE, bold=True)

rev_items = [
    ("Pro订阅", "60%", Inches(3.6), ACCENT_CYAN),
    ("RealmHub抽成", "20%", Inches(1.2), GREEN),
    ("企业版", "15%", Inches(0.9), ACCENT_ORANGE),
    ("家庭版", "5%", Inches(0.3), GRAY),
]
bar_y = Inches(4.5)
bar_x = Inches(0.8)
for name, pct, w, color in rev_items:
    add_shape(slide, bar_x, bar_y, w, Inches(0.5), fill_color=color)
    bar_x = bar_x + w + Inches(0.05)

for i, (name, pct, w, color) in enumerate(rev_items):
    x = Inches(0.8 + i * 1.55)
    add_text(slide, x, Inches(5.1), Inches(1.5), Inches(0.3), f"{name} {pct}", font_size=12, color=color)

# RealmHub ecosystem
add_shape(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.6), fill_color=CARD_BG)
add_text(slide, Inches(1.0), Inches(5.7), Inches(5.0), Inches(0.4),
         "RealmHub 生态系统", font_size=18, color=ACCENT_CYAN, bold=True)

hub_steps = [
    ("创作者", "开发 Realm 包"),
    ("发布", "上传到 RealmHub"),
    ("用户", "购买/下载"),
    ("分成", "创作者获得70%"),
]
for i, (title, desc) in enumerate(hub_steps):
    x = Inches(1.0 + i * 3.0)
    add_text(slide, x, Inches(6.15), Inches(2.5), Inches(0.35), title, font_size=16, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(6.5), Inches(2.5), Inches(0.3), desc, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)
    if i < 3:
        add_text(slide, Inches(3.2 + i * 3.0), Inches(6.15), Inches(0.5), Inches(0.35), "→", font_size=20, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 9: Financial Plan
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "财务规划", "3年财务预测")

# Financial table
add_table(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.55),
          ["年份", "收入", "成本", "净利润"],
          [
              ["2026", "¥50万", "¥504万", "-¥454万"],
              ["2027", "¥500万", "¥1,080万", "-¥580万"],
              ["2028", "¥2,000万", "¥2,280万", "-¥280万"],
          ],
          col_widths=[Inches(2.5), Inches(3.0), Inches(3.0), Inches(3.2)])

# Timeline milestones
add_text(slide, Inches(0.8), Inches(3.6), Inches(5.0), Inches(0.4),
         "关键里程碑", font_size=18, color=WHITE, bold=True)

milestones = [
    ("2026 Q1", "Phase 1 完成\n垂直切片", GREEN),
    ("2026 Q4", "商业化启动\nPro版上线", ACCENT_CYAN),
    ("2027 Q2", "A轮融资\n¥50万 MRR", ACCENT_ORANGE),
    ("2028-2029", "盈亏平衡\n持续增长", WHITE),
]
add_shape(slide, Inches(1.0), Inches(4.55), Inches(11.0), Inches(0.04), fill_color=ACCENT_CYAN)
for i, (time, desc, color) in enumerate(milestones):
    x = Inches(1.0 + i * 3.0)
    add_shape(slide, x + Inches(0.5), Inches(4.4), Inches(0.3), Inches(0.3), fill_color=color, shape_type=MSO_SHAPE.OVAL)
    add_text(slide, x, Inches(4.8), Inches(2.5), Inches(0.3), time, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(5.15), Inches(2.5), Inches(0.6), desc, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# Funding plan
add_text(slide, Inches(0.8), Inches(5.9), Inches(5.0), Inches(0.4),
         "融资计划", font_size=18, color=WHITE, bold=True)

add_card(slide, Inches(0.8), Inches(6.3), Inches(5.5), Inches(1.0),
         "种子轮：¥200万（当前）",
         ["用途：研发60% + 团队25% + 市场15%", "里程碑：10K Stars，RealmHub上线"],
         title_color=ACCENT_CYAN)

add_card(slide, Inches(6.8), Inches(6.3), Inches(5.7), Inches(1.0),
         "A轮：¥1,000-1,500万（2027 Q2）",
         ["里程碑：10K DAU，¥50万月经常性收入"],
         title_color=ACCENT_ORANGE)


# ============================================================
# Slide 10: Milestones & Progress
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "里程碑与当前进度", "Phase 1 已完成，进入 Phase 2")

# Roadmap table
add_table(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5),
          ["阶段", "时间", "状态", "关键交付"],
          [
              ["Phase 1", "2026 Q1", "✅ 完成", "Monorepo, Core, Gateway, CLI"],
              ["Phase 2", "2026 Q2", "🔄 进行中", "Web Dashboard, Realm Matrix UI"],
              ["Phase 3", "2026 Q3", "⏳ 待开始", "RealmHub, 官方Realm包"],
              ["Phase 4", "2026 Q4", "⏳ 待开始", "移动端, Pro版订阅"],
              ["Phase 5", "2027", "⏳ 待开始", "生态扩张, 10K+ DAU"],
          ],
          col_widths=[Inches(2.0), Inches(2.0), Inches(2.0), Inches(5.7)])

# Current metrics
add_text(slide, Inches(0.8), Inches(4.4), Inches(5.0), Inches(0.4),
         "当前指标", font_size=18, color=WHITE, bold=True)

current_metrics = [
    ("📦 8", "核心包"),
    ("🧪 51+", "单元测试"),
    ("📄 MIT", "开源协议"),
    ("🏗️ Phase 1", "垂直切片完成"),
]
for i, (num, label) in enumerate(current_metrics):
    x = Inches(0.8 + i * 3.1)
    add_shape(slide, x, Inches(4.9), Inches(2.8), Inches(1.0), fill_color=CARD_BG)
    add_text(slide, x, Inches(4.95), Inches(2.8), Inches(0.55), num, font_size=28, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(5.5), Inches(2.8), Inches(0.3), label, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# 2026 targets
add_text(slide, Inches(0.8), Inches(6.2), Inches(5.0), Inches(0.4),
         "2026年目标", font_size=18, color=WHITE, bold=True)
targets = ["GitHub Stars: 5,000+", "活跃Realm: 10,000+", "付费转化率: 5%", "全年收入: ¥50万"]
add_bullet_list(slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.8),
                targets, font_size=14, color=GRAY, spacing=Pt(2))


# ============================================================
# Slide 11: Team
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "核心团队", "Team")

# Team cards
team = [
    ("创始人/架构师", "系统设计、Summon机制、产品方向", "10年+系统架构，AI Agent领域深耕", ACCENT_CYAN),
    ("后端负责人", "Core + Storage + Ink Gateway", "Node.js/TypeScript专家，Local-first架构", GREEN),
    ("前端负责人", "Dashboard Web UI", "React/Next.js全栈，数据可视化", ACCENT_ORANGE),
    ("AI研究员", "LLM集成、Agent运行时、Summon引擎", "多模型部署，Prompt Engineering", YELLOW),
    ("产品经理", "RealmHub、用户体验、增长", "个人效率工具，开源社区运营", RGBColor(0xBB, 0x86, 0xFC)),
]

for i, (role, resp, bg_info, color) in enumerate(team):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.5 + row * 2.8)
    add_shape(slide, x, y, Inches(3.7), Inches(2.4), fill_color=CARD_BG, line_color=color)
    add_text(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.3), Inches(0.4), role, font_size=18, color=color, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.65), Inches(3.3), Inches(0.3), "职责", font_size=12, color=GRAY)
    add_text(slide, x + Inches(0.2), y + Inches(0.9), Inches(3.3), Inches(0.5), resp, font_size=14, color=WHITE)
    add_text(slide, x + Inches(0.2), y + Inches(1.45), Inches(3.3), Inches(0.3), "背景", font_size=12, color=GRAY)
    add_text(slide, x + Inches(0.2), y + Inches(1.7), Inches(3.3), Inches(0.5), bg_info, font_size=14, color=WHITE)

# Advisors
add_text(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4),
         "顾问团队：AI伦理专家  |  开源社区顾问  |  效率工具KOL", font_size=14, color=GRAY)


# ============================================================
# Slide 12: Exit Strategy
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "退出机制", "清晰的退出路径")

# Potential acquirers table
add_table(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5),
          ["公司", "动机", "可能性"],
          [
              ["Notion", "补足人生治理维度", "高"],
              ["Microsoft", "增强Copilot个人场景", "高"],
              ["字节跳动", "豆包生态补充", "中"],
              ["OpenClaw基金会", "生态整合", "中"],
          ],
          col_widths=[Inches(3.0), Inches(5.7), Inches(3.0)])

# Exit options (4 quadrants)
options = [
    ("战略收购", "Notion, Microsoft, 字节跳动", Inches(0.8), ACCENT_CYAN),
    ("持续运营", "开源+云服务", Inches(6.8), GREEN),
    ("并购整合", "OpenClaw 生态", Inches(0.8), YELLOW),
    ("IPO", "5-7年, $100M ARR", Inches(6.8), ACCENT_ORANGE),
]
for i, (title, desc, x, color) in enumerate(options):
    row = i // 2
    y = Inches(4.2 + row * 1.5)
    add_shape(slide, x, y, Inches(5.5), Inches(1.2), fill_color=CARD_BG, line_color=color)
    add_text(slide, x + Inches(0.3), y + Inches(0.15), Inches(4.9), Inches(0.4), title, font_size=18, color=color, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.6), Inches(4.9), Inches(0.4), desc, font_size=14, color=GRAY)

# Labels
add_text(slide, Inches(5.5), Inches(4.2), Inches(2.0), Inches(0.3), "高可能性 ↑", font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(5.5), Inches(6.6), Inches(2.0), Inches(0.3), "↓ 长期目标", font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 13: Risk & Mitigation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "风险与应对策略", "Risk & Mitigation")

# Top 3 risks
risks = [
    ("🔴 OpenClaw竞争", "风险：生态庞大，可能覆盖我们的功能", "应对：差异化定位（人生治理 vs 任务自动化），生态兼容", RED),
    ("🟡 LLM成本波动", "风险：API价格上涨影响运营成本", "应对：多供应商策略 + 本地模型（Ollama）降级", YELLOW),
    ("🟢 开源商业化", "风险：代码开源难以收费", "应对：云服务增值 + RealmHub商店 + 企业版服务", GREEN),
]

for i, (title, risk, mitigation, color) in enumerate(risks):
    y = Inches(1.5 + i * 1.9)
    add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.7), fill_color=CARD_BG, line_color=color)
    add_text(slide, Inches(1.1), y + Inches(0.15), Inches(11.1), Inches(0.4), title, font_size=20, color=color, bold=True)
    add_text(slide, Inches(1.1), y + Inches(0.6), Inches(11.1), Inches(0.4), risk, font_size=15, color=WHITE)
    add_text(slide, Inches(1.1), y + Inches(1.05), Inches(11.1), Inches(0.4), mitigation, font_size=15, color=ACCENT_CYAN)

# Bottom note
add_text(slide, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.3),
         "开源商业化已被 GitLab、MongoDB 验证可行", font_size=13, color=GRAY)


# ============================================================
# Slide 14: Vision & Call to Action
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "愿景与行动", "让AI成为人生的治理者")

# Big vision quote
add_shape(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.5), fill_color=CARD_BG, line_color=ACCENT_CYAN)
add_text(slide, Inches(1.5), Inches(1.6), Inches(10.3), Inches(1.3),
         '"让每个人都拥有一个\n真正懂你的AI人生伙伴"',
         font_size=30, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Core differentiation table
add_table(slide, Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.5),
          ["维度", "OpenClaw", "OpenOctopus"],
          [
              ["定位", "任务自动化", "人生治理"],
              ["门槛", "高（写配置）", "低（一键召唤）"],
              ["用户", "技术极客", "生活管理者"],
              ["市场", "小", "大10倍"],
          ],
          col_widths=[Inches(2.5), Inches(4.6), Inches(4.6)])

# Call to action buttons
cta = [
    ("🔗 加入社区", "discord.gg/mwNTk8g5fV", ACCENT_CYAN),
    ("⭐ Star项目", "github.com/open-octopus", ACCENT_ORANGE),
    ("📺 观看Demo", "即将发布", GREEN),
    ("🛠️ 参与建设", "贡献Realm模板", RGBColor(0xBB, 0x86, 0xFC)),
]
for i, (label, url, color) in enumerate(cta):
    x = Inches(0.8 + i * 3.1)
    add_shape(slide, x, Inches(6.0), Inches(2.8), Inches(1.0), fill_color=color)
    add_text(slide, x, Inches(6.05), Inches(2.8), Inches(0.5), label, font_size=18, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(6.5), Inches(2.8), Inches(0.4), url, font_size=12, color=BG_DARK, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 15: Thank You & Q&A
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), fill_color=ACCENT_CYAN)

add_text(slide, Inches(0), Inches(1.5), W, Inches(1.0),
         "OpenOctopus", font_size=54, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(2.5), W, Inches(0.5),
         "Realm-native 人生治理 Agent 系统", font_size=22, color=GRAY, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(5.0), Inches(3.3), Inches(3.333), Inches(0.04), fill_color=ACCENT_CYAN)

add_text(slide, Inches(0), Inches(3.8), W, Inches(0.8),
         "谢谢！\nQuestions & Answers", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Contact info
contacts = [
    "📧 hello@openoctopus.org",
    "🐙 github.com/open-octopus",
    "💬 discord.gg/mwNTk8g5fV",
    "🐦 @OpenOctopusAI",
]
for i, contact in enumerate(contacts):
    add_text(slide, Inches(0), Inches(5.2 + i * 0.45), W, Inches(0.4),
             contact, font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(0), Inches(7.42), W, Inches(0.08), fill_color=ACCENT_CYAN)


# ============================================================
# Save
# ============================================================
output_path = os.path.join(os.path.dirname(__file__), "OpenOctopus-BusinessPlan.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
